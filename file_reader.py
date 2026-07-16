import os
import base64
from pypdf import PdfReader
from docx import Document

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

def extract_pdf_text(file_path, max_chars=15000):
    try:
        reader = PdfReader(file_path)
        text = ""
        for page in reader.pages:
            text += page.extract_text() + "\n"
        return text.strip()[:max_chars]
    except Exception as e:
        return f"Error reading PDF: {str(e)}"

def extract_docx_text(file_path, max_chars=15000):
    try:
        doc = Document(file_path)
        text = "\n".join([para.text for para in doc.paragraphs])
        return text.strip()[:max_chars]
    except Exception as e:
        return f"Error reading DOCX: {str(e)}"

def extract_txt_text(file_path, max_chars=15000):
    try:
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            return f.read().strip()[:max_chars]
    except Exception as e:
        return f"Error reading text file: {str(e)}"

def extract_pptx_text(file_path, max_chars=15000):
    """Extract text from PPTX, slide by slide."""
    if not Presentation:
        return "Error: python-pptx not installed"
    try:
        prs = Presentation(file_path)
        parts = [f"[PPTX: {len(prs.slides)} slides]\n"]
        for i, slide in enumerate(prs.slides):
            slide_texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            slide_texts.append(text)
                if shape.has_table:
                    for row in shape.table.rows[:10]:
                        row_text = " | ".join(cell.text.strip() for cell in row.cells)
                        if row_text.strip():
                            slide_texts.append(row_text)
            if slide_texts:
                parts.append(f"\n--- Slide {i+1} ---\n" + "\n".join(slide_texts))
        return "".join(parts).strip()[:max_chars]
    except Exception as e:
        return f"Error reading PPTX: {str(e)}"

def encode_image_base64(file_path):
    try:
        with open(file_path, "rb") as image_file:
            return base64.b64encode(image_file.read()).decode('utf-8')
    except Exception as e:
        return None

def get_file_type(filename):
    ext = filename.lower().split('.')[-1]
    if ext == 'pdf':
        return 'pdf'
    elif ext == 'docx':
        return 'docx'
    elif ext == 'pptx':
        return 'pptx'
    elif ext in ['txt', 'md']:
        return 'text'
    elif ext in ['png', 'jpg', 'jpeg', 'webp']:
        return 'image'
    else:
        return 'unsupported'

def process_file(file_path, filename):
    file_type = get_file_type(filename)

    if file_type == 'pdf':
        return {'type': 'text', 'content': extract_pdf_text(file_path)}
    elif file_type == 'docx':
        return {'type': 'text', 'content': extract_docx_text(file_path)}
    elif file_type == 'pptx':
        return {'type': 'text', 'content': extract_pptx_text(file_path)}
    elif file_type == 'text':
        return {'type': 'text', 'content': extract_txt_text(file_path)}
    elif file_type == 'image':
        base64_img = encode_image_base64(file_path)
        ext = filename.lower().split('.')[-1]
        mime = f"image/{ext if ext != 'jpg' else 'jpeg'}"
        return {'type': 'image', 'content': base64_img, 'mime': mime}
    else:
        return {'type': 'unsupported', 'content': None}
